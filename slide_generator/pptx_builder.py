from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
from PIL import Image


def get_style_dict():
    return {
        "blue":  {"color": RGBColor(54, 95, 145), "font": "Calibri"},
        "green": {"color": RGBColor(34, 177, 76), "font": "Arial"},
        "orange": {"color": RGBColor(237, 125, 49), "font": "Tahoma"},
    }


def create_presentation(slides_json, output_file, style_name="blue", title_image_path=None):
    style = get_style_dict().get(style_name, get_style_dict()["blue"])
    prs = Presentation()
    # this is the first time I'm working with this module, getting stuff aligned is very hard honestly. especially with images.

    blank_layout = prs.slide_layouts[6]      
    content_slide_layout = prs.slide_layouts[1]  

    for i, slide in enumerate(slides_json.get("slides", [])):
        if i == 0:
            # first (title) slide on a blank layout
            sldr = prs.slides.add_slide(blank_layout)

            # add centered title textbox 
            title_left   = int(prs.slide_width  * 0.10)
            title_width  = int(prs.slide_width  * 0.80)
            title_top    = int(prs.slide_height * 0.08)
            title_height = int(prs.slide_height * 0.18)

            title_box = sldr.shapes.add_textbox(title_left, title_top, title_width, title_height)
            tf = title_box.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = tf.paragraphs[0]
            p.text = slide.get("title", "")
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(44)
                run.font.name = style["font"]
                run.font.color.rgb = style["color"]

            # add image below the title, centered
            if title_image_path and os.path.exists(title_image_path):
                img = Image.open(title_image_path)
                img_width, img_height = img.size
                dpi = img.info.get("dpi", (96, 96))[0]  # fallback if missing

                # convet pixel size -> EMUs
                img_width_emu = Inches(img_width / dpi)
                img_height_emu = Inches(img_height / dpi)

                # scale to fit reserved area
                max_width = prs.slide_width * 0.60
                max_height = prs.slide_height * 0.44
                scale = min(max_width / img_width_emu, max_height / img_height_emu, 1)

                pic_width = int(img_width_emu * scale)
                pic_height = int(img_height_emu * scale)

                left = int((prs.slide_width - pic_width) / 2)
                top = int(title_top + title_height + Pt(16))  # small gap below title

                sldr.shapes.add_picture(
                    title_image_path,
                    left,
                    top,
                    width=pic_width,
                    height=pic_height,
                )

        else:
            # slides 2-7
            sldr = prs.slides.add_slide(content_slide_layout)
            sldr.shapes.title.text = slide.get("title", "")

            content_placeholder = None
            for shape in sldr.shapes:
                if shape.is_placeholder and shape.has_text_frame and shape != sldr.shapes.title:
                    content_placeholder = shape
                    break

            if content_placeholder:
                tf = content_placeholder.text_frame
                tf.clear()
                for bullet in slide.get("bullets", []):
                    if bullet.strip():
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.font.size = Pt(18)
                        p.font.name = style["font"]
                        p.font.color.rgb = style["color"]

    prs.save(output_file)
